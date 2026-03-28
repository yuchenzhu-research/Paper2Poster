import os
import shutil
import subprocess
import sys
import tempfile
import traceback
from time import sleep, time
from types import SimpleNamespace

import json_repair
import Levenshtein
from lxml import etree
from pdf2image import convert_from_path
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.text.text import _Paragraph, _Run
from pptx.util import Length, Pt
from rich import print
from tenacity import RetryCallState, retry, stop_after_attempt, wait_fixed, wait_random

IMAGE_EXTENSIONS = {"bmp", "jpg", "jpeg", "pgm", "png", "ppm", "tif", "tiff", "webp"}

BLACK = RGBColor(0, 0, 0)
YELLOW = RGBColor(255, 255, 0)
BLUE = RGBColor(0, 0, 255)
BORDER_LEN = Pt(2)
BORDER_OFFSET = Pt(2)
LABEL_LEN = Pt(24)
FONT_LEN = Pt(20)


def is_image_path(file: str):
    if file.split(".")[-1].lower() in IMAGE_EXTENSIONS:
        return True
    return False


def get_font_pptcstyle(font: dict):
    font = SimpleNamespace(**font)
    return f"Font Style: bold={font.bold}, italic={font.italic}, underline={font.underline}, size={font.size}pt, color={font.color}, font style={font.name}\n"


def get_font_style(font: dict):
    font = SimpleNamespace(**font)
    styles = []
    if font.size:
        styles.append(f"font-size: {font.size}pt")
    if font.color:
        styles.append(f"color: #{font.color}")
    if font.bold:
        styles.append("font-weight: bold")
    if font.italic:
        styles.append("font-style: italic")
    return "; ".join(styles)


def runs_merge(paragraph: _Paragraph):
    runs = paragraph.runs
    if len(runs) == 0:
        runs = [
            _Run(r, paragraph)
            for r in parse_xml(paragraph._element.xml.replace("fld", "r")).r_lst
        ]
    if len(runs) == 1:
        return runs[0]
    if len(runs) == 0:
        return None
    run = max(runs, key=lambda x: len(x.text))
    run.text = paragraph.text

    for r in runs:
        if r != run:
            r._r.getparent().remove(r._r)
    return run


def older_than(filepath, seconds: int = 10, wait: bool = False):
    if not os.path.exists(filepath):
        while wait:
            print("waiting for:", filepath)
            sleep(1)
            if os.path.exists(filepath):
                sleep(seconds)
                return True
        return False
    file_creation_time = os.path.getctime(filepath)
    current_time = time()
    return seconds < (current_time - file_creation_time)


def edit_distance(text1: str, text2: str):
    return 1 - Levenshtein.distance(text1, text2) / max(len(text1), len(text2))


def get_slide_content(doc_json: dict, slide_title: str, slide: dict):
    slide_desc = slide.get("description", "")
    slide_content = f"Slide Purpose: {slide_title}\nSlide Description: {slide_desc}\n"
    for key in slide.get("subsections", []):
        slide_content += "Slide Content Source: "
        for section in doc_json["sections"]:
            subsections = section.get("subsections", [])
            if isinstance(subsections, dict) and len(subsections) == 1:
                subsections = [
                    {"title": k, "content": v} for k, v in subsections.items()
                ]
            for subsection in subsections:
                try:
                    if edit_distance(key, subsection["title"]) > 0.9:
                        slide_content += f"# {key} \n{subsection['content']}\n"
                except:
                    pass
    return slide_content


def tenacity_log(retry_state: RetryCallState):
    print(retry_state)
    traceback.print_tb(retry_state.outcome.exception().__traceback__)


def get_json_from_response(raw_response: str):
    response = raw_response.strip()
    l, r = response.rfind("```json"), response.rfind("```")
    try:
        if l == -1 or r == -1:
            response = json_repair.loads(response)
        else:
            response = json_repair.loads(response[l + 7 : r].strip())
        return response
    except Exception as e:
        raise RuntimeError("Failed to parse JSON from response", e)


tenacity = retry(
    wait=wait_random(3), stop=stop_after_attempt(5), after=tenacity_log, reraise=True
)


def resolve_soffice_binary() -> str:
    env_path = os.environ.get("SOFFICE_BIN")
    if env_path:
        resolved_env_path = os.path.expanduser(env_path)
        if os.path.isfile(resolved_env_path) and os.access(resolved_env_path, os.X_OK):
            return resolved_env_path
        raise RuntimeError(
            "SOFFICE_BIN is set but does not point to an executable file: "
            f"{resolved_env_path}"
        )

    path_binary = shutil.which("soffice")
    if path_binary:
        return path_binary

    if sys.platform == "darwin":
        macos_app_binary = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if os.path.isfile(macos_app_binary) and os.access(macos_app_binary, os.X_OK):
            return macos_app_binary

    raise RuntimeError(
        "LibreOffice executable not found. Set SOFFICE_BIN, add `soffice` to "
        "PATH, or install LibreOffice. On macOS, the default app bundle path "
        "is /Applications/LibreOffice.app/Contents/MacOS/soffice."
    )


@tenacity
def ppt_to_images(file: str, output_dir: str, warning: bool = False, dpi=72, output_type='png'):
    assert pexists(file), f"File {file} does not exist"
    if pexists(output_dir) and warning:
        print(f"ppt2images: {output_dir} already exists")
    os.makedirs(output_dir, exist_ok=True)
    soffice_bin = resolve_soffice_binary()
    with tempfile.TemporaryDirectory() as temp_dir:
        # Create unique user installation directory for LibreOffice to avoid concurrency issues
        with tempfile.TemporaryDirectory() as user_install_dir:
            command_list = [
                soffice_bin,
                "--headless",
                "--norestore",
                "--nolockcheck",
                f"-env:UserInstallation=file://{user_install_dir}",
                "--convert-to",
                "pdf",
                file,
                "--outdir",
                temp_dir,
            ]
            # Set environment to ensure UTF-8 encoding
            env = os.environ.copy()
            env['LC_ALL'] = 'en_US.UTF-8'
            env['LANG'] = 'en_US.UTF-8'
            subprocess.run(command_list, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, env=env)

        for f in os.listdir(temp_dir):
            if not f.endswith(".pdf"):
                continue
            temp_pdf = pjoin(temp_dir, f)
            images = convert_from_path(temp_pdf, dpi=72)
            for i, img in enumerate(images):
                if output_type == 'png':
                    img.save(pjoin(output_dir, f"poster.png"), 'PNG')
                else:
                    img.save(pjoin(output_dir, f"poster.jpg"), 'JPEG')
            return

        raise RuntimeError("No PDF file was created in the temporary directory", file)


@tenacity
def wmf_to_images(blob: bytes, filepath: str):
    if not filepath.endswith(".jpg"):
        raise ValueError("filepath must end with .jpg")
    dirname = os.path.dirname(filepath)
    basename = os.path.basename(filepath).removesuffix(".jpg")
    soffice_bin = resolve_soffice_binary()
    with tempfile.TemporaryDirectory() as temp_dir:
        with open(pjoin(temp_dir, f"{basename}.wmf"), "wb") as f:
            f.write(blob)
        # Create unique user installation directory for LibreOffice to avoid concurrency issues
        with tempfile.TemporaryDirectory() as user_install_dir:
            command_list = [
                soffice_bin,
                "--headless",
                "--norestore",
                "--nolockcheck",
                f"-env:UserInstallation=file://{user_install_dir}",
                "--convert-to",
                "jpg",
                pjoin(temp_dir, f"{basename}.wmf"),
                "--outdir",
                dirname,
            ]
            env = os.environ.copy()
            env["LC_ALL"] = "en_US.UTF-8"
            env["LANG"] = "en_US.UTF-8"
            subprocess.run(
                command_list,
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                env=env,
            )

    assert pexists(filepath), f"File {filepath} does not exist"


def extract_fill(shape: BaseShape):
    if "fill" not in dir(shape):
        return None
    else:
        return shape.fill._xPr.xml


def apply_fill(shape: BaseShape, fill_xml: str):
    if fill_xml is None:
        return
    new_element = etree.fromstring(fill_xml)
    shape.fill._xPr.getparent().replace(shape.fill._xPr, new_element)


def parse_groupshape(groupshape: GroupShape):
    assert isinstance(groupshape, GroupShape)
    group_top_left_x = groupshape.left
    group_top_left_y = groupshape.top
    group_width = groupshape.width
    group_height = groupshape.height
    shape_top_left_x = min([sp.left for sp in groupshape.shapes])
    shape_top_left_y = min([sp.top for sp in groupshape.shapes])
    shape_width = (
        max([sp.left + sp.width for sp in groupshape.shapes]) - shape_top_left_x
    )
    shape_height = (
        max([sp.top + sp.height for sp in groupshape.shapes]) - shape_top_left_y
    )
    group_shape_xy = []
    for sp in groupshape.shapes:
        group_shape_left = (
            sp.left - shape_top_left_x
        ) * group_width / shape_width + group_top_left_x
        group_shape_top = (
            sp.top - shape_top_left_y
        ) * group_height / shape_height + group_top_left_y
        group_shape_width = sp.width * group_width / shape_width
        group_shape_height = sp.height * group_height / shape_height
        group_shape_xy.append(
            {
                "left": Length(group_shape_left),
                "top": Length(group_shape_top),
                "width": Length(group_shape_width),
                "height": Length(group_shape_height),
            }
        )
    return group_shape_xy


def is_primitive(obj):
    if isinstance(obj, (list, tuple, set, frozenset)):
        return all(is_primitive(item) for item in obj)
    return isinstance(
        obj, (int, float, complex, bool, str, bytes, bytearray, type(None))
    )


DEFAULT_EXCLUDE = set(["element", "language_id", "ln", "placeholder_format"])


def object_to_dict(obj, result=None, exclude=None):
    if result is None:
        result = {}
    exclude = DEFAULT_EXCLUDE.union(exclude or set())
    for attr in dir(obj):
        if attr in exclude:
            continue
        try:
            if not attr.startswith("_") and not callable(getattr(obj, attr)):
                attr_value = getattr(obj, attr)
                if "real" in dir(attr_value):
                    attr_value = attr_value.real
                if attr == "size" and isinstance(attr_value, int):
                    attr_value = Length(attr_value).pt

                if is_primitive(attr_value):
                    result[attr] = attr_value
        except:
            pass
    return result


def merge_dict(d1: dict, d2: list[dict]):
    if len(d2) == 0:
        return d1
    for key in list(d1.keys()):
        values = [d[key] for d in d2]
        if d1[key] is not None and len(values) != 1:
            values.append(d1[key])
        if values[0] is None or not all(value == values[0] for value in values):
            continue
        d1[key] = values[0]
        for d in d2:
            d[key] = None
    return d1


def dict_to_object(dict: dict, obj: object, exclude=None):
    if exclude is None:
        exclude = set()
    for key, value in dict.items():
        if key not in exclude:
            setattr(obj, key, value)


class Config:

    def __init__(self, rundir=None, session_id=None, debug=True):
        self.DEBUG = debug
        if session_id is not None:
            self.set_session(session_id)
        if rundir is not None:
            self.set_rundir(rundir)

    def set_session(self, session_id):
        self.session_id = session_id
        self.set_rundir(f"./runs/{session_id}")

    def set_rundir(self, rundir: str):
        self.RUN_DIR = rundir
        self.IMAGE_DIR = pjoin(self.RUN_DIR, "images")
        for the_dir in [self.RUN_DIR, self.IMAGE_DIR]:
            os.makedirs(the_dir, exist_ok=True)

    def set_debug(self, debug: bool):
        self.DEBUG = debug

    def remove_rundir(self):
        if pexists(self.RUN_DIR):
            shutil.rmtree(self.RUN_DIR)
        if pexists(self.IMAGE_DIR):
            shutil.rmtree(self.IMAGE_DIR)


pjoin = os.path.join
pexists = os.path.exists
pbasename = os.path.basename

if __name__ == "__main__":
    config = Config()
    print(config)
