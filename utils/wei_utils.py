import re
import io
import contextlib
import traceback
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE, MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from camel.types import ModelPlatformType, ModelType
from camel.configs import ChatGPTConfig, QwenConfig, VLLMConfig, OpenRouterConfig, GeminiConfig
import math
from urllib.parse import quote_from_bytes, quote
from PIL import Image
import os
import copy
import io
from utils.src.utils import ppt_to_images
from pathlib import Path
import asyncio
from utils.pptx_utils import *
from utils.critic_utils import *

def get_agent_config(model_type):
    agent_config = {}
    if model_type == 'qwen':
        agent_config = {
            "model_type": ModelType.DEEPINFRA_QWEN_2_5_72B,
            "model_config": QwenConfig().as_dict(),
            "model_platform": ModelPlatformType.DEEPINFRA,
        }
    elif model_type == 'gemini':
        agent_config = {
            "model_type": ModelType.DEEPINFRA_GEMINI_2_FLASH,
            "model_config": GeminiConfig().as_dict(),
            "model_platform": ModelPlatformType.DEEPINFRA,
            'max_images': 99
        }
    elif model_type == 'phi4':
        agent_config = {
            "model_type": ModelType.DEEPINFRA_PHI_4_MULTIMODAL,
            "model_config": QwenConfig().as_dict(),
            "model_platform": ModelPlatformType.DEEPINFRA,
        }
    elif model_type == 'llama-4-scout-17b-16e-instruct':
        agent_config = {
            'model_type': ModelType.ALIYUN_LLAMA4_SCOUT_17B_16E,
            'model_config': QwenConfig().as_dict(),
            'model_platform': ModelPlatformType.QWEN,
            'max_images': 99
        }
    elif model_type == 'qwen-2.5-vl-72b':
        agent_config = {
            'model_type': ModelType.QWEN_2_5_VL_72B,
            'model_config': QwenConfig().as_dict(),
            'model_platform': ModelPlatformType.QWEN,
            'max_images': 99
        }
    elif model_type == 'gemma':
        agent_config = {
            "model_type": "google/gemma-3-4b-it",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:5555/v1',
            'max_images': 99
        }
    elif model_type == 'llava':
        agent_config = {
            "model_type": "llava-hf/llava-onevision-qwen2-7b-ov-hf",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:8000/v1',
            'max_images': 99
        }
    elif model_type == 'molmo-o':
        agent_config = {
            "model_type": "allenai/Molmo-7B-O-0924",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:8000/v1',
            'max_images': 99
        }
    elif model_type == 'qwen-2-vl-7b':
        agent_config = {
            "model_type": "Qwen/Qwen2-VL-7B-Instruct",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:8000/v1',
            'max_images': 99
        }
    elif model_type == 'vllm_phi4':
        agent_config = {
            "model_type": "microsoft/Phi-4-multimodal-instruct",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:8000/v1',
            'max_images': 99
        }
    elif model_type == 'o3-mini':
        agent_config = {
            "model_type": ModelType.O3_MINI,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
        }
    elif model_type == 'gpt-4.1':
        agent_config = {
            "model_type": ModelType.GPT_4_1,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
        }
    elif model_type == 'gpt-4.1-mini':
        agent_config = {
            "model_type": ModelType.GPT_4_1_MINI,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
        }
    elif model_type == '4o':
        agent_config = {
            "model_type": ModelType.GPT_4O,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
            # "model_name": '4o'
        }
    elif model_type == '4o-mini':
        agent_config = {
            "model_type": ModelType.GPT_4O_MINI,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
        }
    elif model_type == 'o1':
        agent_config = {
            "model_type": ModelType.O1,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
            # "model_name": 'o1'
        }
    elif model_type == 'o3':
        agent_config = {
            "model_type": ModelType.O3,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
        }
    elif model_type == 'gpt-5':
        agent_config = {
            "model_type": ModelType.GPT_5,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
        }
    elif model_type in ('gpt-5-mini', 'gpt-5.4-mini'):
        agent_config = {
            "model_type": model_type,
            "model_config": ChatGPTConfig().as_dict(),
            "model_platform": ModelPlatformType.OPENAI,
        }
    elif model_type == 'vllm_qwen_vl':
        agent_config = {
            "model_type": "Qwen/Qwen2.5-VL-7B-Instruct",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:7000/v1'
        }
    elif model_type == 'vllm_qwen':
        agent_config = {
            "model_type": "Qwen/Qwen2.5-7B-Instruct",
            "model_platform": ModelPlatformType.VLLM,
            "model_config": VLLMConfig().as_dict(),
            "url": 'http://localhost:8000/v1',
        }
    elif model_type == 'openrouter_qwen_72b':
        agent_config = {
            'model_type': ModelType.OPENROUTER_QWEN_2_5_72B,
            'model_platform': ModelPlatformType.OPENROUTER,
            'model_config': OpenRouterConfig().as_dict(),
        }
    elif model_type == 'openrouter_qwen_vl_72b':
        agent_config = {
            'model_type': ModelType.OPENROUTER_QWEN_2_5_VL_72B,
            'model_platform': ModelPlatformType.OPENROUTER,
            'model_config': OpenRouterConfig().as_dict(),
        }
    elif model_type == 'openrouter_qwen_vl_7b':
        agent_config = {
            'model_type': ModelType.OPENROUTER_QWEN_2_5_VL_7B,
            'model_platform': ModelPlatformType.OPENROUTER,
            'model_config': OpenRouterConfig().as_dict(),
        }
    elif model_type == 'openrouter_qwen_7b':
        agent_config = {
            'model_type': ModelType.OPENROUTER_QWEN_2_5_7B,
            'model_platform': ModelPlatformType.OPENROUTER,
            'model_config': OpenRouterConfig().as_dict(),
        }
    else:
        agent_config = {
            'model_type': model_type,
            'model_platform': ModelPlatformType.OPENAI_COMPATIBLE_MODEL,
            'model_config': None
        }
    
    return agent_config


def match_response(response):
    response_text = response.msgs[0].content

    # This regular expression looks for text between ```python ... ```
    pattern = r'```python(.*?)```'
    match = re.search(pattern, response_text, flags=re.DOTALL)

    if not match:
        pattern = r'```(.*?)```'
        match = re.search(pattern, response_text, flags=re.DOTALL)

    if match:
        code_snippet = match.group(1).strip()
    else:
        # If there's no fenced code block, fallback to entire response or handle error
        code_snippet = response_text
    return code_snippet

def run_code_with_utils(code, utils_functions):
    return run_code(utils_functions + '\n' + code)

def run_code(code):
    """
    Execute Python code and capture stdout as well as the full stack trace on error.
    Forces __name__ = "__main__" so that if __name__ == "__main__": blocks will run.
    
    Returns:
        (output, error)
        - output: string containing everything that was printed to stdout
        - error: string containing the full traceback if an exception occurred; None otherwise
    """
    stdout_capture = io.StringIO()
    # Provide a globals dict specifying that __name__ is "__main__"
    exec_globals = {"__name__": "__main__"}

    with contextlib.redirect_stdout(stdout_capture):
        try:
            exec(code, exec_globals)
            error = None
        except Exception:
            # Capture the entire stack trace
            error = traceback.format_exc()

    output = stdout_capture.getvalue()
    return output, error


def run_code_from_agent(agent, msg, num_retries=1):
    agent.reset()
    log = []
    for attempt in range(num_retries + 1):  # +1 to include the initial attempt
        response = agent.step(msg)
        code = match_response(response)
        output, error = run_code(code)
        log.append((code, output, error))
        
        if error is None:
            return log
        
        if attempt < num_retries:
            print(f"Retrying... Attempt {attempt + 1} of {num_retries}")
            msg = error
    
    return log

def run_modular(all_code, file_name, with_border=True, with_label=True):
    concatenated_code = utils_functions
    concatenated_code += "\n".join(all_code.values())
    if with_border and with_label:
        concatenated_code += add_border_label_function
        concatenated_code += create_id_map_function
        concatenated_code += save_helper_info_border_label.format(file_name, file_name, file_name)
    elif with_border:
        concatenated_code += add_border_function
        concatenated_code += save_helper_info_border.format(file_name, file_name)
    else:
        concatenated_code += f'\nposter.save("{file_name}")'
    output, error = run_code(concatenated_code)
    return concatenated_code, output, error

def edit_modular(
        agent,
        edit_section_name, 
        feedback,
        all_code, 
        file_name, 
        outline,
        content,
        images,
        actor_prompt,
        num_retries=1,
        prompt_type='initial'
    ):
    agent.reset()
    log = []
    if prompt_type == 'initial':
        msg = actor_prompt.format(
            outline['meta'],
            {edit_section_name: outline[edit_section_name]}, 
            content, 
            images,
            documentation
        )
    elif prompt_type == 'edit':
        assert (edit_section_name == list(feedback.keys())[0])
        msg = actor_prompt.format(
            edit_section_name,
            all_code[edit_section_name],
            feedback,
            {edit_section_name: outline[edit_section_name]}, 
            content, 
            images,
            documentation
        )
    elif prompt_type == 'new':
        assert (list(feedback.keys())[0] == 'all_good')
        msg = actor_prompt.format(
            {edit_section_name: outline[edit_section_name]}, 
            content, 
            images,
            documentation
        )

    for attempt in range(num_retries + 1):
        response = agent.step(msg)
        new_code = match_response(response)
        all_code_changed = all_code.copy()
        all_code_changed[edit_section_name] = new_code
        concatenated_code, output, error = run_modular(all_code_changed, file_name, False, False)
        log.append({
            "code": new_code,
            "output": output,
            "error": error,
            "concatenated_code": concatenated_code
        })
        if error is None:
            return log
        
        if attempt < num_retries:
            print(f"Retrying... Attempt {attempt + 1} of {num_retries}")
            msg = error
            msg += '\nFix your code and try again. The poster is a single-page pptx.'
            if prompt_type != 'initial':
                msg += '\nAssume that you have had a Presentation object named "poster" and a slide named "slide".'

    return log

def add_border_to_all_elements(prs, border_color=RGBColor(255, 0, 0), border_width=Pt(2)):
    """
    Iterates over all slides and shapes in the Presentation object 'prs'
    and applies a red border with the specified width to each shape.
    
    Args:
        prs: The Presentation object to modify.
        border_color: An instance of RGBColor for the border color (default is red).
        border_width: The width of the border as a Pt value (default is 2 points).
    """
    for slide in prs.slides:
        for shape in slide.shapes:
            # Some shapes (like charts or group shapes) might not support border styling
            try:
                # Set the line fill to be solid and assign the desired color and width.
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = border_color
                shape.line.width = border_width
            except Exception as e:
                # If a shape doesn't support setting a border, print a message and continue.
                print(f"Could not add border to shape {shape.shape_type}: {e}")


# 1 point = 12700 EMUs (helper function)
def pt_to_emu(points: float) -> int:
    return int(points * 12700)

def add_border_and_labels(
    prs,
    border_color=RGBColor(255, 0, 0),   # Red border for shapes
    border_width=Pt(2),                # 2-point border width
    label_outline_color=RGBColor(0, 0, 255),  # Blue outline for label circle
    label_text_color=RGBColor(0, 0, 255),     # Blue text color
    label_diameter_pt=40                       # Diameter of the label circle in points
):
    """
    Iterates over all slides and shapes in the Presentation 'prs', applies a 
    red border to each shape, and places a transparent (no fill), blue-outlined 
    circular label with a blue number in the center of each shape. Labels start 
    from 0 and increment for every shape that gets a border.

    Args:
        prs: The Presentation object to modify.
        border_color: RGBColor for the shape border color (default: red).
        border_width: The width of the shape border (Pt).
        label_outline_color: The outline color for the label circle (default: blue).
        label_text_color: The color of the label text (default: blue).
        label_diameter_pt: The diameter of the label circle, in points (default: 40).
    """
    label_diameter_emu = pt_to_emu(label_diameter_pt)  # convert diameter (points) to EMUs
    label_counter = 0  # Start labeling at 0
    labeled_elements = {}

    for slide in prs.slides:
        for shape in slide.shapes:
            # Skip shapes that are labels themselves
            if shape.name.startswith("Label_"):
                continue

            try:
                # --- 1) Add red border to the shape (if supported) ---
                shape.line.fill.solid()
                shape.line.fill.fore_color.rgb = border_color
                shape.line.width = border_width

                # --- 2) Calculate center for the label circle ---
                label_left = shape.left + (shape.width // 2) - (label_diameter_emu // 2)
                label_top  = shape.top  + (shape.height // 2) - (label_diameter_emu // 2)

                # --- 3) Create label circle (an OVAL) in the center of the shape ---
                label_shape = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.OVAL,
                    label_left,
                    label_top,
                    label_diameter_emu,
                    label_diameter_emu
                )
                label_shape.name = f"Label_{label_counter}"  # so we can skip it later

                # **Make the circle completely transparent** (no fill at all)
                label_shape.fill.background()

                # **Give it a blue outline**
                label_shape.line.fill.solid()
                label_shape.line.fill.fore_color.rgb = label_outline_color
                label_shape.line.width = Pt(3)

                # --- 4) Add the label number (centered, blue text) ---
                tf = label_shape.text_frame
                tf.text = str(label_counter)
                paragraph = tf.paragraphs[0]
                paragraph.alignment = PP_ALIGN.CENTER

                run = paragraph.runs[0]
                font = run.font
                font.size = Pt(40)      # Larger font
                font.bold = True
                font.name = "Arial"
                font._element.get_or_change_to_solidFill()
                font.fill.fore_color.rgb = label_text_color
                # Record properties from the original shape and label text.
                labeled_elements[label_counter] = {
                    'left': f'{shape.left} EMU',
                    'top': f'{shape.top} EMU',
                    'width': f'{shape.width} EMU',
                    'height': f'{shape.height} EMU',
                    'font_size': f'{shape.text_frame.font.size} PT' if hasattr(shape, 'text_frame') else None,
                }

                # --- 5) Increment label counter (so every shape has a unique label) ---
                label_counter += 1

            except Exception as e:
                # If the shape doesn't support borders or text, skip gracefully
                print(f"Could not add border/label to shape (type={shape.shape_type}): {e}")

    return labeled_elements


def fill_content(agent, prompt, num_retries, existing_code=''):
    if existing_code == '':
        existing_code = utils_functions
    agent.reset()
    log = []
    cumulative_input_token, cumulative_output_token = 0, 0
    for attempt in range(num_retries + 1):
        response = agent.step(prompt)
        input_token, output_token = account_token(response)
        cumulative_input_token += input_token
        cumulative_output_token += output_token
        new_code = match_response(response)
        all_code = existing_code + '\n' + new_code

        output, error = run_code(all_code)
        log.append({
            "code": new_code,
            "output": output,
            "error": error,
            "concatenated_code": all_code,
            'cumulative_tokens': (cumulative_input_token, cumulative_output_token)
        })

        if error is None:
            return log
        
        if attempt < num_retries:
            print(f"Retrying... Attempt {attempt + 1} of {num_retries}")
            prompt = error
    return log

def apply_theme(agent, prompt, num_retries, existing_code=''):
    return fill_content(agent, prompt, num_retries, existing_code)

def edit_code(agent, prompt, num_retries, existing_code=''):
    return fill_content(agent, prompt, num_retries, existing_code)

def stylize(agent, prompt, num_retries, existing_code=''):
    return fill_content(agent, prompt, num_retries, existing_code)

def gen_layout(agent, prompt, num_retries, name_to_hierarchy, visual_identifier='', existing_code=''):
    if existing_code == '':
        existing_code = utils_functions
    agent.reset()
    log = []
    cumulative_input_token, cumulative_output_token = 0, 0
    for attempt in range(num_retries + 1):
        response = agent.step(prompt)
        input_token, output_token = account_token(response)
        cumulative_input_token += input_token
        cumulative_output_token += output_token
        new_code = match_response(response)
        all_code = existing_code + '\n' + new_code

        # Save visualizations
        all_code += f'''
name_to_hierarchy = {name_to_hierarchy}
identifier = "{visual_identifier}"
get_visual_cues(name_to_hierarchy, identifier)
'''

        output, error = run_code(all_code)
        log.append({
            "code": new_code,
            "output": output,
            "error": error,
            "concatenated_code": all_code,
            'num_tokens': (input_token, output_token),
            'cumulative_tokens': (cumulative_input_token, cumulative_output_token)
        })

        if error is None:
            return log
        
        if attempt < num_retries:
            print(f"Retrying... Attempt {attempt + 1} of {num_retries}")
            prompt = error
    return log

def gen_layout_parallel(agent, prompt, num_retries, existing_code='', slide_width=0, slide_height=0, tmp_name='tmp'):
    if existing_code == '':
        existing_code = utils_functions
        
    existing_code += f'''
poster = create_poster(width_inch={slide_width}, height_inch={slide_height})
slide = add_blank_slide(poster)
save_presentation(poster, file_name="poster_{tmp_name}.pptx")
'''
    agent.reset()
    log = []
    cumulative_input_token, cumulative_output_token = 0, 0
    for attempt in range(num_retries + 1):
        response = agent.step(prompt)
        input_token, output_token = account_token(response)
        cumulative_input_token += input_token
        cumulative_output_token += output_token
        new_code = match_response(response)
        all_code = existing_code + '\n' + new_code

        output, error = run_code(all_code)
        log.append({
            "code": new_code,
            "output": output,
            "error": error,
            "concatenated_code": all_code,
            'num_tokens': (input_token, output_token),
            'cumulative_tokens': (cumulative_input_token, cumulative_output_token)
        })
        if output is None or output == '':
            prompt = 'No object name printed.'
            continue

        if error is None:
            return log
        
        if attempt < num_retries:
            # print(f"Retrying... Attempt {attempt + 1} of {num_retries}", flush=True)
            prompt = error
    return log

def compute_bullet_length(textbox_content):
    total = 0
    for bullet in textbox_content:
        for run in bullet['runs']:
            total += len(run['text'])
    return total

def check_bounding_boxes(bboxes, overall_width, overall_height):
    """
    Given a dictionary 'bboxes' whose keys are bounding-box names and whose values are
    dictionaries with keys 'left', 'top', 'width', and 'height' (all floats),
    along with the overall canvas width and height, this function checks for:

      1) An overlap between any two bounding boxes (it returns a tuple of their names).
      2) A bounding box that extends beyond the overall width or height (it returns a tuple
         containing just that bounding box's name).

    It stops upon finding the first error:
      - If an overlap is found first, it returns (name1, name2).
      - Otherwise, if an overflow is found, it returns (name,).
      - If nothing is wrong, it returns ().

    Parameters:
        bboxes (dict): e.g. {
            "box1": {"left": 10.0, "top": 10.0, "width": 50.0, "height": 20.0},
            "box2": {"left": 55.0, "top": 15.0, "width": 10.0, "height": 10.0},
            ...
        }
        overall_width (float): The total width of the available space.
        overall_height (float): The total height of the available space.

    Returns:
        tuple: Either (box1, box2) if an overlap is found,
               (box,) if a bounding box overflows,
               or () if no problem is found.
    """

    # Convert bboxes into a list of (name, left, top, width, height) for easier iteration.
    box_list = []
    for name, coords in bboxes.items():
        left = coords["left"]
        top = coords["top"]
        width = coords["width"]
        height = coords["height"]
        box_list.append((name, left, top, width, height))

    # Helper function to check overlap between two boxes
    def boxes_overlap(box_a, box_b):
        # Unpack bounding-box data
        name_a, left_a, top_a, width_a, height_a = box_a
        name_b, left_b, top_b, width_b, height_b = box_b

        # Compute right and bottom coordinates
        right_a = left_a + width_a
        bottom_a = top_a + height_a
        right_b = left_b + width_b
        bottom_b = top_b + height_b

        # Rectangles overlap if not separated along either x or y axis
        # If one box is completely to the left or right or above or below the other,
        # there's no overlap.
        no_overlap = (right_a <= left_b or  # A is completely left of B
                      right_b <= left_a or  # B is completely left of A
                      bottom_a <= top_b or  # A is completely above B
                      bottom_b <= top_a)    # B is completely above A
        return not no_overlap

    # 1) Check for overlap first
    n = len(box_list)
    for i in range(n):
        for j in range(i + 1, n):
            if boxes_overlap(box_list[i], box_list[j]):
                return (box_list[i][0], box_list[j][0])  # Return names

    # 2) Check for overflow
    for name, left, top, width, height in box_list:
        right = left + width
        bottom = top + height

        # If boundary is outside [0, overall_width] or [0, overall_height], it's an overflow
        if (left < 0 or top < 0 or right > overall_width or bottom > overall_height):
            return (name,)

    # 3) If nothing is wrong, return empty tuple
    return ()


def is_poster_filled(
    bounding_boxes: dict,
    overall_width: float,
    overall_height: float,
    max_lr_margin: float,
    max_tb_margin: float
) -> bool:
    """
    Given a dictionary of bounding boxes (keys are box names and
    values are dicts with float keys: "left", "top", "width", "height"),
    along with the overall dimensions of the poster and maximum allowed
    margins, this function determines whether the boxes collectively
    fill the poster within those margin constraints.

    :param bounding_boxes: Dictionary of bounding boxes of the form:
                          {
                              "box1": {"left": float, "top": float, "width": float, "height": float},
                              "box2": {...},
                              ...
                          }
    :param overall_width: Total width of the poster
    :param overall_height: Total height of the poster
    :param max_lr_margin: Maximum allowed left and right margins
    :param max_tb_margin: Maximum allowed top and bottom margins
    :return: True if the bounding boxes fill the poster (with no big leftover spaces),
             False otherwise.
    """

    # If there are no bounding boxes, we consider the poster unfilled.
    if not bounding_boxes:
        return False

    # Extract the minimum left, maximum right, minimum top, and maximum bottom from all bounding boxes.
    min_left = min(b["left"] for b in bounding_boxes.values())
    max_right = max(b["left"] + b["width"] for b in bounding_boxes.values())
    min_top = min(b["top"] for b in bounding_boxes.values())
    max_bottom = max(b["top"] + b["height"] for b in bounding_boxes.values())

    # Calculate leftover margins.
    leftover_left = min_left
    leftover_right = overall_width - max_right
    leftover_top = min_top
    leftover_bottom = overall_height - max_bottom

    # Check if leftover margins exceed the allowed maxima.
    if (leftover_left > max_lr_margin or leftover_right > max_lr_margin or
        leftover_top > max_tb_margin or leftover_bottom > max_tb_margin):
        return False

    return True

def check_and_fix_subsections(section, subsections):
    """
    Given a 'section' bounding box and a dictionary of 'subsections',
    checks:

    1) That each subsection is within the main section and that
       no two subsections overlap.
       - If there is a problem, returns a tuple of the names of
         the offending subsections.

    2) That the subsections fully occupy the area of 'section'.
       - If not, greedily expand each subsection (in the order
         left->right->top->bottom), and return a dictionary of
         the updated bounding boxes for the subsections.

    3) Otherwise, returns an empty tuple if nothing is wrong.

    :param section: dict with keys "left", "top", "width", "height".
    :param subsections: dict mapping name -> dict with "left", "top", "width", "height".
    :return: Either
        - tuple of subsection names that are out of bounds or overlapping,
        - dict of expanded bounding boxes if they do not fully occupy 'section',
        - or an empty tuple if everything is correct.
    """

    # --- Utility functions ---
    def right(rect):
        return rect["left"] + rect["width"]

    def bottom(rect):
        return rect["top"] + rect["height"]

    def is_overlapping(r1, r2):
        """
        Returns True if rectangles r1 and r2 overlap (strictly),
        False otherwise.
        """
        return not (
            right(r1) <= r2["left"]
            or r1["left"] >= right(r2)
            or bottom(r1) <= r2["top"]
            or r1["top"] >= bottom(r2)
        )

    # 1) Check each subsection is within the main section
    names_violating = set()
    sec_left, sec_top = section["left"], section["top"]
    sec_right = section["left"] + section["width"]
    sec_bottom = section["top"] + section["height"]

    for name, sub in subsections.items():
        # Check boundary
        sub_left, sub_top = sub["left"], sub["top"]
        sub_right, sub_bottom = right(sub), bottom(sub)
        if (
            sub_left < sec_left
            or sub_top < sec_top
            or sub_right > sec_right
            or sub_bottom > sec_bottom
        ):
            # Out of bounds
            names_violating.add(name)

    # 2) Check pairwise overlaps
    sub_keys = list(subsections.keys())
    for i in range(len(sub_keys)):
        for j in range(i + 1, len(sub_keys)):
            n1, n2 = sub_keys[i], sub_keys[j]
            if is_overlapping(subsections[n1], subsections[n2]):
                # Mark both as violating
                names_violating.add(n1)
                names_violating.add(n2)

    # If anything violated boundaries or overlapped, return them as a tuple
    if names_violating:
        return tuple(sorted(names_violating))

    # 3) Check if subsections fully occupy the section by area.
    #    (Since we've checked there's no overlap, area-based check is safe for "full coverage".)
    area_section = section["width"] * section["height"]
    area_subs = sum(
        sub["width"] * sub["height"] for sub in subsections.values()
    )

    if area_subs < area_section:
        # -- We need to expand subsections greedily. --

        # Make a copy of the bounding boxes so as not to modify originals.
        expanded_subs = {
            name: {
                "left": sub["left"],
                "top": sub["top"],
                "width": sub["width"],
                "height": sub["height"],
            }
            for name, sub in subsections.items()
        }

        # Helper to see whether we are touching a boundary or another subsection
        def touching_left(sname, sbox):
            if abs(sbox["left"] - sec_left) < 1e-9:
                # touches main section left boundary
                return True
            # touches the right edge of another subsection
            for oname, obox in expanded_subs.items():
                if oname == sname:
                    continue
                if abs(right(obox) - sbox["left"]) < 1e-9:
                    return True
            return False

        def touching_right(sname, sbox):
            r = right(sbox)
            if abs(r - sec_right) < 1e-9:
                return True
            for oname, obox in expanded_subs.items():
                if oname == sname:
                    continue
                if abs(obox["left"] - r) < 1e-9:
                    return True
            return False

        def touching_top(sname, sbox):
            if abs(sbox["top"] - sec_top) < 1e-9:
                return True
            for oname, obox in expanded_subs.items():
                if oname == sname:
                    continue
                if abs(bottom(obox) - sbox["top"]) < 1e-9:
                    return True
            return False

        def touching_bottom(sname, sbox):
            b = bottom(sbox)
            if abs(b - sec_bottom) < 1e-9:
                return True
            for oname, obox in expanded_subs.items():
                if oname == sname:
                    continue
                if abs(obox["top"] - b) < 1e-9:
                    return True
            return False

        # Attempt a single pass of expansions, left->right->top->bottom
        for name in expanded_subs:
            sub = expanded_subs[name]

            # Expand left if not touching left boundary or another box
            if not touching_left(name, sub):
                # The "left boundary" is the maximum "right" of any subsection strictly to the left,
                # or the section's left boundary, whichever is larger.
                left_bound = sec_left
                for oname, obox in expanded_subs.items():
                    if oname == name:
                        continue
                    r_ = obox["left"] + obox["width"]
                    # only consider those that are strictly left of this sub
                    if r_ <= sub["left"] and r_ > left_bound:
                        left_bound = r_
                # Now expand
                delta = sub["left"] - left_bound
                if delta > 1e-9:  # If there's any real gap
                    sub["width"] += delta
                    sub["left"] = left_bound

            # Expand right if not touching right boundary or another box
            if not touching_right(name, sub):
                right_bound = sec_right
                sub_right = sub["left"] + sub["width"]
                for oname, obox in expanded_subs.items():
                    if oname == name:
                        continue
                    left_ = obox["left"]
                    # only consider those that are strictly to the right
                    if left_ >= sub_right and left_ < right_bound:
                        right_bound = left_
                delta = right_bound - (sub["left"] + sub["width"])
                if delta > 1e-9:
                    sub["width"] += delta

            # Expand top if not touching top boundary or another box
            if not touching_top(name, sub):
                top_bound = sec_top
                for oname, obox in expanded_subs.items():
                    if oname == name:
                        continue
                    b_ = obox["top"] + obox["height"]
                    if b_ <= sub["top"] and b_ > top_bound:
                        top_bound = b_
                delta = sub["top"] - top_bound
                if delta > 1e-9:
                    sub["height"] += delta
                    sub["top"] = top_bound

            # Expand bottom if not touching bottom boundary or another box
            if not touching_bottom(name, sub):
                bottom_bound = sec_bottom
                sub_bottom = sub["top"] + sub["height"]
                for oname, obox in expanded_subs.items():
                    if oname == name:
                        continue
                    other_top = obox["top"]
                    if other_top >= sub_bottom and other_top < bottom_bound:
                        bottom_bound = other_top
                delta = bottom_bound - (sub["top"] + sub["height"])
                if delta > 1e-9:
                    sub["height"] += delta

        # After expansion, return the expanded dictionary
        # per the spec: "If the second case happens, return a dictionary ...
        # containing the modified bounding box dictionaries."
        return expanded_subs

    # If we get here, then area_subs == area_section and there's no overlap => all good
    return ()

async def rendered_dims(html: Path) -> tuple[int, int]:
    from playwright.async_api import async_playwright

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page    = await browser.new_page()        # no fixed viewport yet
        resolved = html.resolve()
        # quote_from_bytes expects bytes, so we encode the path as UTF‐8:
        url = "file://" + quote_from_bytes(str(resolved).encode("utf-8"), safe="/:")
        await page.goto(url, wait_until="networkidle")

        # 1) bounding-box of <body>
        body_box = await page.eval_on_selector(
            "body",
            "el => el.getBoundingClientRect()")
        w = int(body_box["width"])
        h = int(body_box["height"])

        await browser.close()
        return w, h

    
def html_to_png(html_abs_path, poster_width_default, poster_height_default, output_path):
    from playwright.sync_api import sync_playwright

    html_file = html_abs_path

    try:
        w, h = asyncio.run(rendered_dims(html_file))
    except:
        w = poster_width_default
        h = poster_height_default

    with sync_playwright() as p:
        path_posix = Path(html_file).resolve().as_posix()

        file_url = "file://" + quote(path_posix, safe="/:")
        browser = p.chromium.launch()
        page    = browser.new_page(viewport={"width": w, "height": h})
        page.goto(file_url, wait_until='networkidle')
        page.screenshot(path=output_path, full_page=True)
        browser.close()

def account_token(response):
    input_token = response.info['usage']['prompt_tokens']
    output_token = response.info['usage']['completion_tokens']

    return input_token, output_token

def style_bullet_content(bullet_content_item, color, fill_color):
    for i in range(len(bullet_content_item)):
        bullet_content_item[i]['runs'][0]['color'] = color
        bullet_content_item[i]['runs'][0]['fill_color'] = fill_color

def scale_to_target_area(width, height, target_width=900, target_height=1200):
    """
    Scale the given width and height by the same factor to achieve a new area equal 
    to target_width * target_height while preserving the aspect ratio.

    Parameters:
      width (float or int): The original width.
      height (float or int): The original height.
      target_width (int, optional): The target width for area calculation. Default is 900.
      target_height (int, optional): The target height for area calculation. Default is 1200.

    Returns:
      tuple: (new_width, new_height) after scaling such that the area is target_width * target_height.
    """
    # Calculate target area from provided dimensions.
    target_area = target_width * target_height
    
    # Calculate original area
    current_area = width * height
    
    # Compute scale factor required: s^2 * (width * height) = target_area => s = sqrt(target_area / (width * height))
    scale_factor = math.sqrt(target_area / current_area)
    
    # Calculate new dimensions
    new_width = width * scale_factor
    new_height = height * scale_factor
    
    # Optional: Round the dimensions to integers.
    return int(round(new_width)), int(round(new_height))

def char_capacity(
    bbox,
    font_size_px=40 * (96 / 72),  # Default font size in px (40pt converted to px)
    *,
    # Average glyph width as fraction of font-size (≈0.6 for monospace,
    # ≈0.52–0.55 for most proportional sans-serif faces)
    avg_width_ratio: float = 0.54,
    line_height_ratio: float = 1,
    # Optional inner padding in px that the renderer might reserve
    padding_px: int = 0,
) -> int:
    """
    Estimate the number of characters that will fit into a rectangular text box.

    Parameters
    ----------
    bbox : (x, y, height, width)  # all in pixels
    font_size_px : int           # font size in px
    avg_width_ratio : float      # average char width ÷ fontSize
    line_height_ratio : float    # line height ÷ fontSize
    padding_px : int             # optional inner padding on each side

    Returns
    -------
    int : estimated character capacity
    """
    CHAR_CONST = 10
    _, _, height_px, width_px = bbox

    usable_w = max(0, width_px - 2 * padding_px)
    usable_h = max(0, height_px - 2 * padding_px)

    if usable_w == 0 or usable_h == 0:
        return 0  # box is too small

    avg_char_w = font_size_px * avg_width_ratio
    line_height = font_size_px * line_height_ratio

    chars_per_line = max(1, math.floor(usable_w / avg_char_w))
    lines = max(1, math.floor(usable_h / line_height))

    return chars_per_line * lines * CHAR_CONST

def estimate_characters(width_in_inches, height_in_inches, font_size_points, line_spacing_points=None):
    """
    Estimate the number of characters that can fit into a bounding box.

    :param width_in_inches:  The width of the bounding box, in inches.
    :param height_in_inches: The height of the bounding box, in inches.
    :param font_size_points: The font size, in points.
    :param line_spacing_points: (Optional) The line spacing, in points.
                                Defaults to 1.5 × font_size_points if not provided.
    :return: Estimated number of characters that fit in the bounding box.
    """
    if line_spacing_points is None:
        # Default line spacing is 1.5 times the font size
        line_spacing_points = 1.5 * font_size_points

    # 1 inch = 72 points 
    width_in_points = width_in_inches * 72
    height_in_points = height_in_inches * 72

    # Rough approximation of the average width of a character: half of the font size
    avg_char_width = 0.5 * font_size_points

    # Number of characters that can fit per line
    chars_per_line = int(width_in_points // avg_char_width)

    # Number of lines that can fit in the bounding box
    lines_count = int(height_in_points // line_spacing_points)

    # Total number of characters
    total_characters = chars_per_line * lines_count

    return total_characters

def equivalent_length_with_forced_breaks(text, width_in_inches, font_size_points):
    """
    Returns the "width-equivalent length" of the text when forced newlines
    are respected. Each physical line (including partial) is counted as if it
    had 'max_chars_per_line' characters.
    
    This number can exceed len(text), because forced newlines waste leftover
    space on the line.
    """
    # 1 inch = 72 points
    width_in_points = width_in_inches * 72
    avg_char_width = 0.5 * font_size_points

    # How many characters fit in one fully occupied line?
    max_chars_per_line = int(width_in_points // avg_char_width)

    # Split on explicit newlines
    logical_lines = text.split('\n')

    total_equiv_length = 0

    for line in logical_lines:
        # If the line is empty, we still "use" one line (which is max_chars_per_line slots).
        if not line:
            total_equiv_length += max_chars_per_line
            continue

        line_length = len(line)
        # How many sub-lines (wraps) does it need?
        sub_lines = math.ceil(line_length / max_chars_per_line)

        # Each sub-line is effectively counted as if it were fully used
        total_equiv_length += sub_lines * max_chars_per_line

    return total_equiv_length

def actual_rendered_length(
    text,
    width_in_inches,
    height_in_inches,
    font_size_points,
    line_spacing_points=None
):
    """
    Estimate how many characters from `text` will actually fit in the bounding
    box, taking into account explicit newlines.
    """
    if line_spacing_points is None:
        line_spacing_points = 1.5 * font_size_points

    # 1 inch = 72 points
    width_in_points = width_in_inches * 72
    height_in_points = height_in_inches * 72

    # Estimate average character width
    avg_char_width = 0.5 * font_size_points

    # Maximum chars per line (approx)
    max_chars_per_line = int(width_in_points // avg_char_width)

    # Maximum number of lines that can fit
    max_lines = int(height_in_points // line_spacing_points)

    # Split on newline chars to get individual "logical" lines
    logical_lines = text.split('\n')

    used_lines = 0
    displayed_chars = 0

    for line in logical_lines:
        # If the line is empty, it still takes one printed line
        if not line:
            used_lines += 1
            # Stop if we exceed available lines
            if used_lines >= max_lines:
                break
            continue

        # Number of sub-lines the text will occupy if it wraps
        sub_lines = math.ceil(len(line) / max_chars_per_line)

        # If we don't exceed the bounding box's vertical capacity
        if used_lines + sub_lines <= max_lines:
            # All chars fit within the bounding box
            displayed_chars += len(line)
            used_lines += sub_lines
        else:
            # Only part of this line will fit
            lines_left = max_lines - used_lines
            if lines_left <= 0:
                # No space left at all
                break

            # We can render only `lines_left` sub-lines of this line
            # That means we can render up to:
            chars_that_fit = lines_left * max_chars_per_line

            # Clip to the actual number of characters
            chars_that_fit = min(chars_that_fit, len(line))

            displayed_chars += chars_that_fit
            used_lines += lines_left  # We've used up all remaining lines
            break  # No more space in the bounding box

    return displayed_chars


def remove_hierarchy_and_id(data):
    """
    Recursively remove the 'hierarchy' and 'id' fields from a nested
    dictionary representing sections and subsections.
    """
    if isinstance(data, dict):
        # Create a new dict to store filtered data
        new_data = {}
        for key, value in data.items():
            # Skip the keys "hierarchy" and "id"
            if key in ("hierarchy", "id", 'location'):
                continue
            # Recursively process the value
            new_data[key] = remove_hierarchy_and_id(value)
        return new_data
    elif isinstance(data, list):
        # If it's a list, process each item recursively
        return [remove_hierarchy_and_id(item) for item in data]
    else:
        # Base case: if it's neither dict nor list, just return the value as is
        return data
    
def outline_estimate_num_chars(outline):
    for k, v in outline.items():
        if k == 'meta':
            continue
        if 'title' in k.lower() or 'author' in k.lower() or 'reference' in k.lower():
            continue
        if not 'subsections' in v:
            num_chars = estimate_characters(
                v['location']['width'], 
                v['location']['height'], 
                60, line_spacing_points=None
            )
            v['num_chars'] = num_chars
        else:
            for k_sub, v_sub in v['subsections'].items():
                if 'title' in k_sub.lower():
                    continue
                if 'path' in v_sub:
                    continue
                num_chars = estimate_characters(
                    v_sub['location']['width'], 
                    v_sub['location']['height'], 
                    60, line_spacing_points=None
                )
                v_sub['num_chars'] = num_chars

def generate_length_suggestions(result_json, original_section_outline, raw_section_outline):
    NOT_CHANGE = 'Do not change text.'
    original_section_outline = json.loads(original_section_outline)
    suggestion_flag = False
    new_section_outline = copy.deepcopy(result_json)
    def check_length(text, target, width, height):
        text_length = equivalent_length_with_forced_breaks(
            text,
            width,
            font_size_points=60,
        )
        if text_length - target > 100:
            return f'Text too long, shrink by {text_length - target} characters.'
        elif target - text_length > 100:
            return f'Text too short, expand by {target - text_length} characters.'
        else:
            return NOT_CHANGE

    if 'num_chars' in original_section_outline:
        new_section_outline['suggestions'] = check_length(
            result_json['description'], 
            original_section_outline['num_chars'],
            raw_section_outline['location']['width'],
            raw_section_outline['location']['height']
        )
        if new_section_outline['suggestions'] != NOT_CHANGE:
            suggestion_flag = True
    if 'subsections' in original_section_outline:
        for k, v in original_section_outline['subsections'].items():
            if 'num_chars' in v:
                new_section_outline['subsections'][k]['suggestion'] = check_length(
                    result_json['subsections'][k]['description'], 
                    v['num_chars'],
                    raw_section_outline['subsections'][k]['location']['width'],
                    raw_section_outline['subsections'][k]['location']['height']
                )
                if new_section_outline['subsections'][k]['suggestion'] != NOT_CHANGE:
                    suggestion_flag = True

    return new_section_outline, suggestion_flag

def get_img_ratio(img_path):
    img = Image.open(img_path)
    return {
        'width': img.width,
        'height': img.height
    }

def get_img_ratio_in_section(content_json):
    res = {}
    if 'path' in content_json:
        res[content_json['path']] = get_img_ratio(content_json['path'])

    if 'subsections' in content_json:
        for subsection_name, val in content_json['subsections'].items():
            if 'path' in val:
                res[val['path']] = get_img_ratio(val['path'])

    return res


def get_snapshot_from_section(leaf_section, section_name, name_to_hierarchy, leaf_name, section_code, empty_poster_path='poster.pptx'):
    hierarchy = name_to_hierarchy[leaf_name]
    hierarchy_overflow_name = f'tmp/overflow_check_<{section_name}>_<{leaf_section}>_hierarchy_{hierarchy}'
    run_code_with_utils(section_code, utils_functions)
    poster = Presentation(empty_poster_path)
    # add border regardless of the hierarchy
    curr_location = add_border_hierarchy(
        poster, 
        name_to_hierarchy, 
        hierarchy, 
        border_width=10,
        # regardless=True
    )
    if not leaf_section in curr_location:
        leaf_section = section_name
    save_presentation(poster, file_name=f"{hierarchy_overflow_name}.pptx")
    ppt_to_images(
        f"{hierarchy_overflow_name}.pptx", 
        hierarchy_overflow_name, 
        dpi=200
    )
    poster_image_path = os.path.join(f"{hierarchy_overflow_name}", "slide_0001.jpg")
    poster_image = Image.open(poster_image_path)

    poster_width = emu_to_inches(poster.slide_width)
    poster_height = emu_to_inches(poster.slide_height)
    locations = convert_pptx_bboxes_json_to_image_json(
        curr_location, 
        poster_width, 
        poster_height
    )
    zoomed_in_img = zoom_in_image_by_bbox(
        poster_image, 
        locations[leaf_name], 
        padding=0.01
    )
    # save the zoomed_in_img
    zoomed_in_img.save(f"{hierarchy_overflow_name}_zoomed_in.jpg")
    return curr_location, zoomed_in_img, f"{hierarchy_overflow_name}_zoomed_in.jpg"
